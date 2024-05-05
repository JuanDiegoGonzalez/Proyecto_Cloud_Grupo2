import os
from io import BytesIO
import time
from gcloud import storage
from google.cloud import pubsub_v1
from celery import Celery
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from odf import text, teletype
from odf.opendocument import load
from sqlalchemy import create_engine
from sqlalchemy import text as tx

os.environ.setdefault("GCLOUD_PROJECT", "proyectocloud")
os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = os.path.join(os.getcwd(), "backend", "file_processing", 'storagesa.json')  # TODO: Verificar que el archivo exista en /backend/views/

DATABASE_URL = 'postgresql://postgres:password@34.31.55.58:5432/postgres'  # BDURL

def crear_pdf():
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size = 12)
    return pdf

def actualizar_bd(id_tarea):
    engine = create_engine(DATABASE_URL)
    engine.connect()
    query = f"UPDATE tarea SET status = 'PROCESSED' WHERE id = {id_tarea}"
    with engine.begin() as conn:
        conn.execute(tx(query))

def docx_a_pdf(docx_file, pdf_file, id_tarea):
    gcs = storage.Client()
    bucket = gcs.get_bucket("bucketproyectocloud")
    blob_download = bucket.blob(docx_file)

    document = Document(BytesIO(blob_download.download_as_string()))
    pdf = crear_pdf()
    for para in document.paragraphs:
        text = para.text.encode('latin-1', 'replace').decode('latin-1')
        pdf.cell(200, 10, txt=text, ln=True, align='L')
    
    pdf_bytes = pdf.output(dest='S').encode('latin-1')

    blob_upload = bucket.blob(pdf_file)
    blob_upload.upload_from_string(pdf_bytes, content_type='application/pdf')
    actualizar_bd(id_tarea)

def pptx_a_pdf(pptx_file, pdf_file, id_tarea):
    gcs = storage.Client()
    bucket = gcs.get_bucket("bucketproyectocloud")
    blob_download = bucket.blob(pptx_file)

    presentation = Presentation(BytesIO(blob_download.download_as_string()))
    pdf = crear_pdf()
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.cell(200, 10, txt=shape.text, ln=True, align='L')

    pdf_bytes = pdf.output(dest='S').encode('latin-1')

    blob_upload = bucket.blob(pdf_file)
    blob_upload.upload_from_string(pdf_bytes, content_type='application/pdf')
    actualizar_bd(id_tarea)

def xlsx_a_pdf(xlsx_file, pdf_file, id_tarea):
    gcs = storage.Client()
    bucket = gcs.get_bucket("bucketproyectocloud")
    blob_download = bucket.blob(xlsx_file)

    wb = load_workbook(BytesIO(blob_download.download_as_string()))
    pdf = crear_pdf()
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        for row in sheet.iter_rows():
            for cell in row:
                pdf.cell(40, 10, txt=str(cell.value), ln=True)

    pdf_bytes = pdf.output(dest='S').encode('latin-1')

    blob_upload = bucket.blob(pdf_file)
    blob_upload.upload_from_string(pdf_bytes, content_type='application/pdf')
    actualizar_bd(id_tarea)

def odt_a_pdf(odt_file, pdf_file, id_tarea):
    gcs = storage.Client()
    bucket = gcs.get_bucket("bucketproyectocloud")
    blob_download = bucket.blob(odt_file)

    doc = load(BytesIO(blob_download.download_as_string()))
    pdf = crear_pdf()
    for para in doc.getElementsByType(text.P):
        pp = teletype.extractText(para).encode('latin-1', 'replace').decode('latin-1')
        pdf.cell(200, 10, txt=pp, ln=True, align='L')
    
    pdf_bytes = pdf.output(dest='S').encode('latin-1')

    blob_upload = bucket.blob(pdf_file)
    blob_upload.upload_from_string(pdf_bytes, content_type='application/pdf')
    actualizar_bd(id_tarea)


def callback(message):
    print(f"Received message: {message}")
    message.ack()

    parts = message.data.decode("utf-8").split(";")

    match parts[0]:
        case "docx":
            docx_a_pdf(parts[1], parts[2], parts[3])

        case "pptx":
            pptx_a_pdf(parts[1], parts[2], parts[3])

        case "xlsx":
            xlsx_a_pdf(parts[1], parts[2], parts[3])
        
        case "odt":
            odt_a_pdf(parts[1], parts[2], parts[3])

        case _:
            ...


project_id = "proyectocloud-422409"
subscription_name = "proyecto3pubsub-sub"

subscriber = pubsub_v1.SubscriberClient()
subscription_path = subscriber.subscription_path(project_id, subscription_name)

subscriber.subscribe(subscription_path, callback=callback)

print("Listening for messages on subscription...")
while True:
    time.sleep(2)  # Add a delay to reduce CPU usage
