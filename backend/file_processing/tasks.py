import os
from io import BytesIO
from gcloud import storage
from celery import Celery
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from odf import text, teletype
from odf.opendocument import load
from sqlalchemy import create_engine
from sqlalchemy import text as tx

os.environ.setdefault("GCLOUD_PROJECT", "entrega3cloud")
os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = '/readings/backend/views/storagesa.json'  # TODO: Verificar que el archivo exista en /backend/views/

app = Celery('tasks', broker = 'redis://redis:6379/0')  # TODO: Poner url de la mv?
DATABASE_URL = 'postgresql://postgres:password@35.232.145.254:5432/postgres'  # BDURL

def crear_pdf():
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size = 12)
    return pdf

def actualizar_bd(id_tarea):
    engine = create_engine(DATABASE_URL)
    engine.connect()
    query = f"UPDATE tarea SET status = 'PROCESSED' WHERE id = {id_tarea}"
    with engine.begin() as conn:
        conn.execute(tx(query))

@app.task
def docx_a_pdf(docx_file, pdf_file, id_tarea):
    gcs = storage.Client()
    bucket = gcs.get_bucket("bucketproyecto3cloud")
    blob_download = bucket.blob(docx_file)

    document = Document(BytesIO(blob_download.download_as_string()))
    pdf = crear_pdf()
    for para in document.paragraphs:
        text = para.text.encode('latin-1', 'replace').decode('latin-1')
        pdf.cell(200, 10, txt=text, ln=True, align='L')
    
    pdf_bytes = pdf.output(dest='S').encode('latin-1')

    blob_upload = bucket.blob(pdf_file)
    blob_upload.upload_from_string(pdf_bytes, content_type='application/pdf')
    actualizar_bd(id_tarea)

@app.task
def pptx_a_pdf(pptx_file, pdf_file, id_tarea):
    gcs = storage.Client()
    bucket = gcs.get_bucket("bucketproyecto3cloud")
    blob_download = bucket.blob(pptx_file)

    presentation = Presentation(BytesIO(blob_download.download_as_string()))
    pdf = crear_pdf()
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.cell(200, 10, txt=shape.text, ln=True, align='L')

    pdf_bytes = pdf.output(dest='S').encode('latin-1')

    blob_upload = bucket.blob(pdf_file)
    blob_upload.upload_from_string(pdf_bytes, content_type='application/pdf')
    actualizar_bd(id_tarea)

@app.task
def xlsx_a_pdf(xlsx_file, pdf_file, id_tarea):
    gcs = storage.Client()
    bucket = gcs.get_bucket("bucketproyecto3cloud")
    blob_download = bucket.blob(xlsx_file)

    wb = load_workbook(BytesIO(blob_download.download_as_string()))
    pdf = crear_pdf()
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        for row in sheet.iter_rows():
            for cell in row:
                pdf.cell(40, 10, txt=str(cell.value), ln=True)

    pdf_bytes = pdf.output(dest='S').encode('latin-1')

    blob_upload = bucket.blob(pdf_file)
    blob_upload.upload_from_string(pdf_bytes, content_type='application/pdf')
    actualizar_bd(id_tarea)

@app.task
def odt_a_pdf(odt_file, pdf_file, id_tarea):
    gcs = storage.Client()
    bucket = gcs.get_bucket("bucketproyecto3cloud")
    blob_download = bucket.blob(odt_file)

    doc = load(BytesIO(blob_download.download_as_string()))
    pdf = crear_pdf()
    for para in doc.getElementsByType(text.P):
        pp = teletype.extractText(para).encode('latin-1', 'replace').decode('latin-1')
        pdf.cell(200, 10, txt=pp, ln=True, align='L')
    
    pdf_bytes = pdf.output(dest='S').encode('latin-1')

    blob_upload = bucket.blob(pdf_file)
    blob_upload.upload_from_string(pdf_bytes, content_type='application/pdf')
    actualizar_bd(id_tarea)
